"""
Generador de presentación PowerPoint – AgroSmart Andino
Grupo 2 | Curso Big Data | Maestría IA – UNI
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta de colores ────────────────────────────────────────
VERDE_OSCURO  = RGBColor(0x1B, 0x5E, 0x20)   # #1B5E20
VERDE_MEDIO   = RGBColor(0x2E, 0x7D, 0x32)   # #2E7D32
VERDE_CLARO   = RGBColor(0x4C, 0xAF, 0x50)   # #4CAF50
VERDE_FONDO   = RGBColor(0xE8, 0xF5, 0xE9)   # #E8F5E9
AMARILLO      = RGBColor(0xFF, 0xC1, 0x07)   # #FFC107
NARANJA       = RGBColor(0xFF, 0x57, 0x22)   # #FF5722
BLANCO        = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_OSCURO   = RGBColor(0x21, 0x21, 0x21)
GRIS_CLARO    = RGBColor(0xF5, 0xF5, 0xF5)
AZUL          = RGBColor(0x01, 0x57, 0x9B)   # #01579B

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completamente en blanco

# ── Helpers ──────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=GRIS_OSCURO,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=16, title=None, title_size=18,
                   bullet_color=GRIS_OSCURO, title_color=VERDE_OSCURO,
                   bullet_char="▸ "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
    for b in bullets:
        para = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = bullet_char + b
        run.font.size = Pt(font_size)
        run.font.color.rgb = bullet_color
    return txBox

def header_bar(slide, title, subtitle=None,
               bar_color=VERDE_OSCURO, title_color=BLANCO, sub_color=AMARILLO):
    """Barra superior con título"""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), fill_color=bar_color)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.08), Inches(11), Inches(0.6),
                font_size=28, bold=True, color=title_color, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.65), Inches(11), Inches(0.4),
                    font_size=14, bold=False, color=sub_color, align=PP_ALIGN.LEFT)

def footer_bar(slide, text="AgroSmart Andino | Grupo 2 | Big Data – Maestría IA UNI | Mayo 2026"):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), fill_color=VERDE_MEDIO)
    add_textbox(slide, text,
                Inches(0.3), Inches(7.17), Inches(12.5), Inches(0.3),
                font_size=10, color=BLANCO, align=PP_ALIGN.CENTER)

def slide_number(slide, n):
    add_textbox(slide, str(n),
                Inches(12.8), Inches(7.15), Inches(0.4), Inches(0.3),
                font_size=10, color=BLANCO, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════
# SLIDE 1 – PORTADA
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Fondo degradado simulado con dos rectángulos
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=VERDE_OSCURO)
add_rect(sl, 0, Inches(4.5), SLIDE_W, Inches(3), fill_color=VERDE_MEDIO)

# Línea decorativa
add_rect(sl, Inches(0.5), Inches(1.5), Inches(12.3), Pt(3), fill_color=AMARILLO)

# Título principal
add_textbox(sl, "🌱 AgroSmart Andino",
            Inches(0.5), Inches(0.2), Inches(12), Inches(1.2),
            font_size=44, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)

# Subtítulo
add_textbox(sl,
    "Pipeline Big Data con Dashboard en Tiempo Real\n"
    "para el Monitoreo Inteligente de Papa Nativa\n"
    "en Huangascar, Yauyos – 3,200 m.s.n.m.",
    Inches(0.5), Inches(1.6), Inches(12), Inches(1.8),
    font_size=20, bold=False, color=VERDE_FONDO, align=PP_ALIGN.LEFT)

# Separador
add_rect(sl, Inches(0.5), Inches(3.5), Inches(12.3), Pt(2), fill_color=AMARILLO)

# Curso / docente
add_textbox(sl,
    "Curso: Big Data   |   Docente: Mg. Rosa Virginia Encinas Quille",
    Inches(0.5), Inches(3.6), Inches(12), Inches(0.4),
    font_size=14, color=AMARILLO, align=PP_ALIGN.LEFT)

# Integrantes
integrantes = [
    "Montes Espinoza, Anahys        – MIA-2024-001",
    "Vizcardo Flores, Christian      – MIA-2024-002",
    "Massa Quispe, Cristhian         – MIA-2024-003",
    "Huali Condori, Freddy           – MIA-2024-004",
    "Lazaro Galarza, Wilmer          – MIA-2024-005",
]
add_bullet_box(sl, integrantes,
               Inches(0.5), Inches(4.1), Inches(9), Inches(2.8),
               font_size=14, bullet_color=BLANCO,
               bullet_char="  · ")

add_textbox(sl, "Grupo 2   |   Lima, Perú — Mayo 2026",
            Inches(0.5), Inches(7.0), Inches(10), Inches(0.4),
            font_size=13, color=VERDE_FONDO, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=GRIS_CLARO)
header_bar(sl, "Agenda", "Estructura de la presentación — 15 minutos")
footer_bar(sl); slide_number(sl, 2)

agenda = [
    ("1", "El Problema", "Contexto agrícola andino y necesidad"),
    ("2", "Propuesta de Valor", "Solución innovadora AgroSmart Andino"),
    ("3", "Requerimientos", "Funcionales y no funcionales"),
    ("4", "Arquitectura", "Pipeline Big Data de extremo a extremo"),
    ("5", "Datos y ETL", "Fuentes: NASA POWER, IoT, MIDAGRI — Spark en EMR"),
    ("6", "Modelos ML", "GBT Regressor: R²=0.6764, RMSE=0.97 t/ha"),
    ("7", "Cassandra NoSQL", "Capa de servicio: 3,794 registros, baja latencia"),
    ("8", "Dashboard", "Streamlit: alertas, riego, predicciones en tiempo real"),
    ("9", "Resultados", "Métricas consolidadas del sistema"),
    ("10", "Viabilidad", "Escalabilidad, impacto y expansión futura"),
    ("11", "Conclusiones", "Experiencias del equipo y aprendizajes"),
]

for i, (num, titulo, desc) in enumerate(agenda):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    lft = Inches(0.4 + col * 6.5)
    top = Inches(1.3 + row * 0.87)
    add_rect(sl, lft, top, Inches(0.5), Inches(0.55), fill_color=VERDE_OSCURO)
    add_textbox(sl, num, lft, top + Pt(6), Inches(0.5), Inches(0.45),
                font_size=14, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_textbox(sl, titulo, lft + Inches(0.55), top, Inches(5.5), Inches(0.3),
                font_size=14, bold=True, color=VERDE_OSCURO)
    add_textbox(sl, desc, lft + Inches(0.55), top + Inches(0.3), Inches(5.5), Inches(0.28),
                font_size=11, color=GRIS_OSCURO)

# ════════════════════════════════════════════════════════════
# SLIDE 3 – EL PROBLEMA
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=GRIS_CLARO)
header_bar(sl, "El Problema", "¿Qué desafío enfrenta la agricultura andina?")
footer_bar(sl); slide_number(sl, 3)

# Caja contexto
add_rect(sl, Inches(0.3), Inches(1.2), Inches(5.8), Inches(5.7), fill_color=BLANCO,
         line_color=VERDE_CLARO, line_width=1.5)
add_textbox(sl, "Contexto: Huangascar, Yauyos",
            Inches(0.4), Inches(1.3), Inches(5.5), Inches(0.45),
            font_size=16, bold=True, color=VERDE_OSCURO)
add_bullet_box(sl, [
    "3,200 m.s.n.m. — zona papicultora andina",
    "5 variedades de papa nativa (Solanum tuberosum)",
    "~500 familias agricultoras dependientes",
    "Temporada agrícola: julio – mayo",
    "Acceso limitado a tecnología y datos",
], Inches(0.4), Inches(1.75), Inches(5.5), Inches(2.5),
   font_size=14, bullet_color=GRIS_OSCURO)

# Stats de impacto
stats = [
    ("38%", "reducción de cosecha\npor helada no anticipada (2022)"),
    ("0", "sistemas de alerta\ntemprana disponibles localmente"),
    ("3°C", "temperatura mínima crítica\npara daño en tuberización"),
]
for i, (val, desc) in enumerate(stats):
    top = Inches(1.75 + i * 1.15)
    add_rect(sl, Inches(0.4), top, Inches(1.2), Inches(1.0), fill_color=NARANJA)
    add_textbox(sl, val, Inches(0.4), top + Inches(0.05), Inches(1.2), Inches(0.6),
                font_size=22, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, Inches(1.65), top + Inches(0.1), Inches(4.1), Inches(0.8),
                font_size=12, color=GRIS_OSCURO)

# Caja amenazas
add_rect(sl, Inches(6.4), Inches(1.2), Inches(6.5), Inches(5.7), fill_color=BLANCO,
         line_color=NARANJA, line_width=1.5)
add_textbox(sl, "⚠ Principales amenazas",
            Inches(6.5), Inches(1.3), Inches(6.2), Inches(0.45),
            font_size=16, bold=True, color=NARANJA)
amenazas = [
    ("🌨 Heladas", "T < 3°C durante tuberización destruye el follaje"),
    ("🍂 Tizón tardío", "Phytophthora infestans: >85% HR + T 10-25°C"),
    ("💧 Déficit hídrico", "ET₀ > precipitación → riego ineficiente"),
    ("📊 Sin datos", "Decisiones basadas solo en experiencia empírica"),
    ("🌡 Cambio climático", "Mayor frecuencia e intensidad de eventos extremos"),
]
for i, (ico, desc) in enumerate(amenazas):
    top = Inches(1.8 + i * 1.0)
    add_textbox(sl, ico, Inches(6.5), top, Inches(1.8), Inches(0.45),
                font_size=14, bold=True, color=VERDE_OSCURO)
    add_textbox(sl, desc, Inches(8.3), top, Inches(4.4), Inches(0.45),
                font_size=13, color=GRIS_OSCURO)

# ════════════════════════════════════════════════════════════
# SLIDE 4 – PROPUESTA DE VALOR
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=VERDE_OSCURO)
header_bar(sl, "Propuesta de Valor", "¿Por qué AgroSmart Andino es diferente?",
           bar_color=VERDE_MEDIO, title_color=BLANCO, sub_color=AMARILLO)
footer_bar(sl); slide_number(sl, 4)

add_textbox(sl, "AgroSmart Andino transforma datos climáticos e IoT en decisiones agrícolas accionables\n"
            "para comunidades andinas, reduciendo pérdidas y optimizando recursos.",
            Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.7),
            font_size=15, color=AMARILLO, italic=True)

# 6 tarjetas de valor
tarjetas = [
    (VERDE_CLARO, "🎯 Alerta temprana",
     "Detección automática de heladas,\ntizón tardío y déficit hídrico\ncon semáforo de 3 niveles"),
    (VERDE_CLARO, "🤖 ML predictivo",
     "GBT Regressor predicción de\nrendimiento por variedad\nR²=0.6764 en EMR real"),
    (VERDE_CLARO, "💧 Riego inteligente",
     "Recomendaciones de riego\ncalculadas con ET₀ (Penman-Monteith)\ny humedad de suelo IoT"),
    (AZUL, "☁ Pipeline Big Data",
     "NASA POWER + IoT → S3 Raw\n→ Spark EMR → S3 Gold\n→ Cassandra → Dashboard"),
    (AZUL, "⚡ Baja latencia",
     "Apache Cassandra 4.1\n3,794 registros consultables\nen < 10 ms"),
    (AZUL, "📊 Dashboard 24/7",
     "Streamlit con 11 componentes\nvisuales: KPIs, mapas de riesgo,\npredicciones y alertas activas"),
]
for i, (col, tit, desc) in enumerate(tarjetas):
    c = i % 3
    r = i // 3
    lft = Inches(0.35 + c * 4.3)
    top = Inches(2.0 + r * 2.35)
    add_rect(sl, lft, top, Inches(4.1), Inches(2.2), fill_color=BLANCO)
    add_rect(sl, lft, top, Inches(4.1), Inches(0.45), fill_color=col)
    add_textbox(sl, tit, lft + Inches(0.1), top + Inches(0.05),
                Inches(3.9), Inches(0.35),
                font_size=14, bold=True, color=BLANCO)
    add_textbox(sl, desc, lft + Inches(0.1), top + Inches(0.5),
                Inches(3.9), Inches(1.6),
                font_size=12, color=GRIS_OSCURO)

# ════════════════════════════════════════════════════════════
# SLIDE 5 – REQUERIMIENTOS
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=GRIS_CLARO)
header_bar(sl, "Requerimientos del Sistema",
           "Funcionales y No Funcionales")
footer_bar(sl); slide_number(sl, 5)

# Funcionales
add_rect(sl, Inches(0.3), Inches(1.2), Inches(6.3), Inches(5.7),
         fill_color=BLANCO, line_color=VERDE_CLARO, line_width=1.5)
add_rect(sl, Inches(0.3), Inches(1.2), Inches(6.3), Inches(0.5), fill_color=VERDE_OSCURO)
add_textbox(sl, "✅  Requerimientos Funcionales",
            Inches(0.4), Inches(1.25), Inches(6.0), Inches(0.4),
            font_size=15, bold=True, color=BLANCO)
rf = [
    "RF-01: Ingesta automática NASA POWER API (10 años)",
    "RF-02: Ingesta telemetría IoT (52,632 reg. horarios)",
    "RF-03: ETL PySpark: normalización, unión, cálculo índices",
    "RF-04: Cálculo IHB, ITT, IEH, ET₀ para papa andina",
    "RF-05: Entrenamiento RF + GBT con PySpark MLlib",
    "RF-06: Predicción de rendimiento por variedad (t/ha)",
    "RF-07: Generación de alertas: HELADA / TIZÓN / RIEGO",
    "RF-08: Almacenamiento NoSQL en Cassandra 4.1",
    "RF-09: Dashboard Streamlit con visualización interactiva",
    "RF-10: Recomendaciones de riego automatizadas",
]
add_bullet_box(sl, rf, Inches(0.4), Inches(1.8), Inches(6.0), Inches(4.8),
               font_size=12.5, bullet_color=GRIS_OSCURO, bullet_char="  · ")

# No Funcionales
add_rect(sl, Inches(6.85), Inches(1.2), Inches(6.1), Inches(5.7),
         fill_color=BLANCO, line_color=AZUL, line_width=1.5)
add_rect(sl, Inches(6.85), Inches(1.2), Inches(6.1), Inches(0.5), fill_color=AZUL)
add_textbox(sl, "⚙  Requerimientos No Funcionales",
            Inches(6.95), Inches(1.25), Inches(5.8), Inches(0.4),
            font_size=15, bold=True, color=BLANCO)
rnf = [
    ("Escalabilidad", "EMR auto-scaling: 1–10 nodos worker"),
    ("Rendimiento", "Procesamiento 72,397 registros en < 5 min"),
    ("Disponibilidad", "Cassandra replication factor = 1 (demo)"),
    ("Seguridad", "Credenciales en .env, S3 IAM policies"),
    ("Tolerancia fallos", "Reintentos PySpark, logs en S3"),
    ("Usabilidad", "Dashboard intuitivo sin capacitación técnica"),
    ("Interoperabilidad", "APIs REST: NASA POWER, MIDAGRI CSV"),
    ("Portabilidad", "Docker Cassandra, Python 3.11, EMR 6.15"),
]
for i, (cat, desc) in enumerate(rnf):
    top = Inches(1.8 + i * 0.59)
    add_textbox(sl, cat + ":", Inches(6.95), top, Inches(1.9), Inches(0.45),
                font_size=12, bold=True, color=AZUL)
    add_textbox(sl, desc, Inches(8.85), top, Inches(3.8), Inches(0.45),
                font_size=12, color=GRIS_OSCURO)

# ════════════════════════════════════════════════════════════
# SLIDE 6 – ARQUITECTURA
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=GRIS_CLARO)
header_bar(sl, "Arquitectura Big Data", "Pipeline de extremo a extremo — 5 capas")
footer_bar(sl); slide_number(sl, 6)

# Capa 1: Ingesta
add_rect(sl, Inches(0.2), Inches(1.25), Inches(2.3), Inches(5.6), fill_color=BLANCO,
         line_color=VERDE_CLARO, line_width=1.5)
add_rect(sl, Inches(0.2), Inches(1.25), Inches(2.3), Inches(0.5), fill_color=VERDE_OSCURO)
add_textbox(sl, "1. INGESTA", Inches(0.25), Inches(1.3), Inches(2.2), Inches(0.4),
            font_size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
fuentes = ["🌍 NASA POWER API\n(10 años, 5 estaciones)",
           "📡 Sensores IoT\n(52,632 reg. horarios)",
           "📋 MIDAGRI CSV\n(1,500 reg. productivos)"]
for i, f in enumerate(fuentes):
    top = Inches(1.85 + i * 1.65)
    add_rect(sl, Inches(0.3), top, Inches(2.1), Inches(1.5), fill_color=VERDE_FONDO)
    add_textbox(sl, f, Inches(0.35), top + Inches(0.1), Inches(2.0), Inches(1.3),
                font_size=11.5, color=VERDE_OSCURO, align=PP_ALIGN.CENTER)

# Flecha
add_textbox(sl, "→", Inches(2.55), Inches(3.8), Inches(0.5), Inches(0.5),
            font_size=24, bold=True, color=AMARILLO, align=PP_ALIGN.CENTER)

# Capa 2: Almacenamiento Raw
add_rect(sl, Inches(3.1), Inches(1.25), Inches(2.0), Inches(5.6), fill_color=BLANCO,
         line_color=AMARILLO, line_width=1.5)
add_rect(sl, Inches(3.1), Inches(1.25), Inches(2.0), Inches(0.5), fill_color=AMARILLO)
add_textbox(sl, "2. S3 DATA LAKE", Inches(3.15), Inches(1.3), Inches(1.9), Inches(0.4),
            font_size=12, bold=True, color=VERDE_OSCURO, align=PP_ALIGN.CENTER)
capas_s3 = [("🗂 RAW\ns3://agrosmart-\nbucket/raw/", GRIS_CLARO),
            ("⚙ PROCESSED\ns3://.../processed/", AMARILLO),
            ("⭐ GOLD\ns3://.../gold/\n(CSV analíticos)", VERDE_CLARO)]
for i, (txt, col) in enumerate(capas_s3):
    top = Inches(1.85 + i * 1.65)
    add_rect(sl, Inches(3.2), top, Inches(1.8), Inches(1.45), fill_color=col)
    add_textbox(sl, txt, Inches(3.25), top + Inches(0.1), Inches(1.7), Inches(1.3),
                font_size=10.5, color=VERDE_OSCURO, align=PP_ALIGN.CENTER)

# Flecha
add_textbox(sl, "→", Inches(5.15), Inches(3.8), Inches(0.5), Inches(0.5),
            font_size=24, bold=True, color=AMARILLO, align=PP_ALIGN.CENTER)

# Capa 3: Procesamiento
add_rect(sl, Inches(5.7), Inches(1.25), Inches(2.3), Inches(5.6), fill_color=BLANCO,
         line_color=NARANJA, line_width=1.5)
add_rect(sl, Inches(5.7), Inches(1.25), Inches(2.3), Inches(0.5), fill_color=NARANJA)
add_textbox(sl, "3. EMR / SPARK", Inches(5.75), Inches(1.3), Inches(2.2), Inches(0.4),
            font_size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
proc = ["🔥 Spark 3.4 en\nEMR 6.15 (AWS)\n3 nodos m5.xlarge",
        "📐 ETL + índices\nIHB, ITT, IEH, ET₀\n72,397 registros",
        "🤖 ML PySpark\nRandomForest + GBT\nR²=0.6764"]
for i, p in enumerate(proc):
    top = Inches(1.85 + i * 1.65)
    add_rect(sl, Inches(5.8), top, Inches(2.1), Inches(1.45), fill_color=RGBColor(0xFF,0xEB,0xD3))
    add_textbox(sl, p, Inches(5.85), top + Inches(0.05), Inches(2.0), Inches(1.35),
                font_size=11, color=GRIS_OSCURO, align=PP_ALIGN.CENTER)

# Flecha
add_textbox(sl, "→", Inches(8.05), Inches(3.8), Inches(0.5), Inches(0.5),
            font_size=24, bold=True, color=AMARILLO, align=PP_ALIGN.CENTER)

# Capa 4: Cassandra
add_rect(sl, Inches(8.6), Inches(1.25), Inches(2.0), Inches(5.6), fill_color=BLANCO,
         line_color=AZUL, line_width=1.5)
add_rect(sl, Inches(8.6), Inches(1.25), Inches(2.0), Inches(0.5), fill_color=AZUL)
add_textbox(sl, "4. CASSANDRA", Inches(8.65), Inches(1.3), Inches(1.9), Inches(0.4),
            font_size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
cass_tabs = ["📊 alertas_riesgo\n(1,825 filas)",
             "🌡 metricas_clima\n(1,964 filas)",
             "🌾 predicciones_\nrendimiento\n(5 variedades)"]
for i, t in enumerate(cass_tabs):
    top = Inches(1.85 + i * 1.65)
    add_rect(sl, Inches(8.7), top, Inches(1.8), Inches(1.45), fill_color=RGBColor(0xE3,0xF2,0xFD))
    add_textbox(sl, t, Inches(8.75), top + Inches(0.05), Inches(1.7), Inches(1.35),
                font_size=10.5, color=AZUL, align=PP_ALIGN.CENTER)

# Flecha
add_textbox(sl, "→", Inches(10.65), Inches(3.8), Inches(0.5), Inches(0.5),
            font_size=24, bold=True, color=AMARILLO, align=PP_ALIGN.CENTER)

# Capa 5: Dashboard
add_rect(sl, Inches(11.2), Inches(1.25), Inches(1.9), Inches(5.6), fill_color=BLANCO,
         line_color=VERDE_CLARO, line_width=1.5)
add_rect(sl, Inches(11.2), Inches(1.25), Inches(1.9), Inches(0.5), fill_color=VERDE_CLARO)
add_textbox(sl, "5. STREAMLIT", Inches(11.25), Inches(1.3), Inches(1.8), Inches(0.4),
            font_size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
dash_comp = ["🚦 Semáforo\nde riesgo",
             "💧 Riego\nautomático",
             "📈 Predicciones\nML en vivo"]
for i, d in enumerate(dash_comp):
    top = Inches(1.85 + i * 1.65)
    add_rect(sl, Inches(11.3), top, Inches(1.7), Inches(1.45), fill_color=VERDE_FONDO)
    add_textbox(sl, d, Inches(11.35), top + Inches(0.15), Inches(1.6), Inches(1.2),
                font_size=11, color=VERDE_OSCURO, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 7 – DATOS Y ETL
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=GRIS_CLARO)
header_bar(sl, "Datos & Pipeline ETL",
           "Fuentes heterogéneas → Gold Layer | Apache Spark 3.4 en Amazon EMR 6.15")
footer_bar(sl); slide_number(sl, 7)

# Tabla de fuentes
headers = ["Fuente", "Tipo", "Registros", "Variables clave"]
rows_data = [
    ["NASA POWER API", "Climático histórico\n(2015–2024)", "18,250 diarios\n5 estaciones", "T_max, T_min, Prec, HR, RS"],
    ["Sensores IoT", "Telemetría horaria\n(sintética calibrada)", "52,632 horarios\n3 nodos IoT", "HR suelo, T suelo, pH, N, P, K"],
    ["MIDAGRI", "Estadísticas agrícolas\n(Yauyos 2015–2024)", "1,500 productivos\n5 variedades", "Rendimiento t/ha, área siembra"],
]
col_widths = [Inches(2.5), Inches(2.5), Inches(2.5), Inches(5.0)]
col_lefts  = [Inches(0.3), Inches(2.85), Inches(5.4), Inches(7.95)]

# Header de tabla
for j, (h, w, l) in enumerate(zip(headers, col_widths, col_lefts)):
    add_rect(sl, l, Inches(1.3), w, Inches(0.45), fill_color=VERDE_OSCURO)
    add_textbox(sl, h, l + Inches(0.05), Inches(1.35), w - Inches(0.1), Inches(0.35),
                font_size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows_data):
    bg = BLANCO if i % 2 == 0 else VERDE_FONDO
    for j, (cell, w, l) in enumerate(zip(row, col_widths, col_lefts)):
        add_rect(sl, l, Inches(1.75 + i * 0.8), w, Inches(0.75), fill_color=bg,
                 line_color=VERDE_CLARO, line_width=0.5)
        add_textbox(sl, cell, l + Inches(0.05), Inches(1.78 + i * 0.8), w - Inches(0.1), Inches(0.7),
                    font_size=11, color=GRIS_OSCURO, align=PP_ALIGN.CENTER)

# Total registros
add_rect(sl, Inches(0.3), Inches(4.2), Inches(12.65), Inches(0.5), fill_color=VERDE_MEDIO)
add_textbox(sl, "TOTAL: 72,382 registros procesados en el pipeline Spark",
            Inches(0.4), Inches(4.25), Inches(12.5), Inches(0.4),
            font_size=14, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# Índices calculados
add_textbox(sl, "Índices Agroclimáticos Calculados con PySpark:",
            Inches(0.3), Inches(4.85), Inches(12.5), Inches(0.4),
            font_size=15, bold=True, color=VERDE_OSCURO)
indices = [
    ("IHB", "Índice de Helada Biológica", "IHB = 1 si T_min < 3°C"),
    ("ITT", "Índice Tizón Tardío", "ITT = f(HR>85%, T 10–25°C)"),
    ("IEH", "Índice Estrés Hídrico", "IEH = Déficit hídrico / ET₀"),
    ("ET₀", "Evapotranspiración ref.", "Penman-Monteith (FAO-56)"),
]
for i, (sig, nombre, formula) in enumerate(indices):
    lft = Inches(0.3 + i * 3.2)
    add_rect(sl, lft, Inches(5.3), Inches(3.1), Inches(1.65), fill_color=BLANCO,
             line_color=VERDE_CLARO, line_width=1.0)
    add_rect(sl, lft, Inches(5.3), Inches(3.1), Inches(0.45), fill_color=VERDE_CLARO)
    add_textbox(sl, sig, lft + Inches(0.05), Inches(5.35), Inches(3.0), Inches(0.35),
                font_size=15, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_textbox(sl, nombre, lft + Inches(0.05), Inches(5.78), Inches(3.0), Inches(0.35),
                font_size=12, bold=True, color=VERDE_OSCURO)
    add_textbox(sl, formula, lft + Inches(0.05), Inches(6.12), Inches(3.0), Inches(0.75),
                font_size=11, color=GRIS_OSCURO, italic=True)

# ════════════════════════════════════════════════════════════
# SLIDE 8 – MODELOS ML
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=GRIS_CLARO)
header_bar(sl, "Modelos de Machine Learning",
           "PySpark MLlib sobre Amazon EMR 6.15 — Ejecución real: 23 mayo 2026")
footer_bar(sl); slide_number(sl, 8)

# Modelo 1: Clasificación
add_rect(sl, Inches(0.3), Inches(1.2), Inches(6.0), Inches(5.75), fill_color=BLANCO,
         line_color=NARANJA, line_width=1.5)
add_rect(sl, Inches(0.3), Inches(1.2), Inches(6.0), Inches(0.5), fill_color=NARANJA)
add_textbox(sl, "MODELO 1: Clasificación de Riesgo",
            Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.4),
            font_size=15, bold=True, color=BLANCO)
add_textbox(sl, "RandomForestClassifier (PySpark MLlib)",
            Inches(0.4), Inches(1.8), Inches(5.8), Inches(0.35),
            font_size=13, bold=True, color=NARANJA)
add_bullet_box(sl, [
    "Target: nivel_riesgo (BAJO / MEDIO / ALTO / CRÍTICO)",
    "Features: IHB, ITT, IEH, T_min, HR, ET₀",
    "numTrees=100, maxDepth=10",
    "Train/Test: 80%/20% — stratified split",
], Inches(0.4), Inches(2.2), Inches(5.8), Inches(2.0),
   font_size=13, bullet_color=GRIS_OSCURO)

# Métricas clasificación
add_textbox(sl, "Métricas:", Inches(0.4), Inches(4.3), Inches(5.8), Inches(0.35),
            font_size=13, bold=True, color=NARANJA)
metricas_clf = [("Accuracy", "~87%"), ("Precision", "~0.86"), ("Recall", "~0.87")]
for i, (m, v) in enumerate(metricas_clf):
    lft = Inches(0.4 + i * 1.85)
    add_rect(sl, lft, Inches(4.7), Inches(1.75), Inches(0.9), fill_color=RGBColor(0xFF,0xEB,0xD3))
    add_textbox(sl, v, lft + Inches(0.05), Inches(4.72), Inches(1.65), Inches(0.45),
                font_size=18, bold=True, color=NARANJA, align=PP_ALIGN.CENTER)
    add_textbox(sl, m, lft + Inches(0.05), Inches(5.18), Inches(1.65), Inches(0.35),
                font_size=11, color=GRIS_OSCURO, align=PP_ALIGN.CENTER)

add_textbox(sl,
    "⚠ Nota: RandomForest en modo demostrativo.\n"
    "Métricas estimadas sobre datos sintéticos calibrados.",
    Inches(0.4), Inches(5.75), Inches(5.8), Inches(0.9),
    font_size=11, color=GRIS_OSCURO, italic=True)

# Modelo 2: Regresión
add_rect(sl, Inches(6.65), Inches(1.2), Inches(6.3), Inches(5.75), fill_color=BLANCO,
         line_color=VERDE_OSCURO, line_width=1.5)
add_rect(sl, Inches(6.65), Inches(1.2), Inches(6.3), Inches(0.5), fill_color=VERDE_OSCURO)
add_textbox(sl, "MODELO 2: Regresión de Rendimiento",
            Inches(6.75), Inches(1.25), Inches(6.1), Inches(0.4),
            font_size=15, bold=True, color=BLANCO)
add_textbox(sl, "GBT Regressor — Gradient Boosted Trees",
            Inches(6.75), Inches(1.8), Inches(6.1), Inches(0.35),
            font_size=13, bold=True, color=VERDE_OSCURO)
add_bullet_box(sl, [
    "Target: rendimiento_tm_ha (toneladas/hectárea)",
    "Features: T_media, T_min, Prec, HR, ET₀, altitud",
    "maxIter=50, maxDepth=5, stepSize=0.1",
    "Train/Test: 80%/20% — random split",
    "Ejecución REAL en AWS EMR cluster 6.15",
], Inches(6.75), Inches(2.2), Inches(6.1), Inches(2.2),
   font_size=13, bullet_color=GRIS_OSCURO)

# Métricas regresión destacadas
metricas_reg = [
    ("R²", "0.6764", "Coef. determinación"),
    ("RMSE", "0.97 t/ha", "Error cuadrático medio"),
    ("MAE", "~0.73 t/ha", "Error absoluto medio"),
]
for i, (nom, val, desc) in enumerate(metricas_reg):
    lft = Inches(6.75 + i * 2.05)
    add_rect(sl, lft, Inches(4.55), Inches(1.95), Inches(1.2), fill_color=VERDE_FONDO,
             line_color=VERDE_CLARO, line_width=1.0)
    add_textbox(sl, val, lft + Inches(0.05), Inches(4.6), Inches(1.85), Inches(0.55),
                font_size=20, bold=True, color=VERDE_OSCURO, align=PP_ALIGN.CENTER)
    add_textbox(sl, nom, lft + Inches(0.05), Inches(5.18), Inches(1.85), Inches(0.3),
                font_size=12, bold=True, color=VERDE_OSCURO, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, lft + Inches(0.05), Inches(5.5), Inches(1.85), Inches(0.2),
                font_size=9, color=GRIS_OSCURO, align=PP_ALIGN.CENTER)

# Predicciones por variedad
add_textbox(sl, "Predicciones por Variedad — Campaña 2025:",
            Inches(6.75), Inches(5.85), Inches(6.1), Inches(0.35),
            font_size=13, bold=True, color=VERDE_OSCURO)
variedades = [
    ("Peruanita", "11.3 t/ha"), ("Huayro", "10.8 t/ha"), ("Canchan", "9.5 t/ha"),
    ("Amarilla", "8.2 t/ha"), ("Nativa Roja", "6.3 t/ha")
]
for i, (var, pred) in enumerate(variedades):
    lft = Inches(6.75 + i * 1.22)
    add_rect(sl, lft, Inches(6.25), Inches(1.15), Inches(0.65), fill_color=VERDE_FONDO)
    add_textbox(sl, pred, lft + Inches(0.03), Inches(6.28), Inches(1.09), Inches(0.3),
                font_size=13, bold=True, color=VERDE_OSCURO, align=PP_ALIGN.CENTER)
    add_textbox(sl, var, lft + Inches(0.03), Inches(6.58), Inches(1.09), Inches(0.28),
                font_size=9, color=GRIS_OSCURO, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 9 – CASSANDRA NoSQL
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=GRIS_CLARO)
header_bar(sl, "Capa de Servicio: Apache Cassandra 4.1",
           "NoSQL de baja latencia para alertas y consultas en tiempo real")
footer_bar(sl); slide_number(sl, 9)

# Izquierda: info del cluster
add_rect(sl, Inches(0.3), Inches(1.2), Inches(5.5), Inches(5.7), fill_color=BLANCO,
         line_color=AZUL, line_width=1.5)
add_rect(sl, Inches(0.3), Inches(1.2), Inches(5.5), Inches(0.5), fill_color=AZUL)
add_textbox(sl, "Configuración del Cluster",
            Inches(0.4), Inches(1.25), Inches(5.3), Inches(0.4),
            font_size=15, bold=True, color=BLANCO)
cluster_info = [
    ("Servidor", "EC2 t2.medium — 54.159.9.135"),
    ("Cassandra", "4.1.9 (Docker compose)"),
    ("Keyspace", "agrosmart"),
    ("Replication", "SimpleStrategy, RF=1"),
    ("CQL", "CREATE TABLE IF NOT EXISTS..."),
]
for i, (k, v) in enumerate(cluster_info):
    top = Inches(1.8 + i * 0.7)
    add_textbox(sl, k + ":", Inches(0.4), top, Inches(1.5), Inches(0.5),
                font_size=13, bold=True, color=AZUL)
    add_textbox(sl, v, Inches(1.95), top, Inches(3.7), Inches(0.5),
                font_size=13, color=GRIS_OSCURO)

add_textbox(sl, "Total registros cargados:",
            Inches(0.4), Inches(5.5), Inches(5.3), Inches(0.4),
            font_size=14, bold=True, color=AZUL)
add_rect(sl, Inches(0.4), Inches(5.95), Inches(2.0), Inches(0.8), fill_color=AZUL)
add_textbox(sl, "3,794", Inches(0.4), Inches(5.98), Inches(2.0), Inches(0.5),
            font_size=24, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
add_textbox(sl, "filas totales en 3 tablas", Inches(2.5), Inches(6.0), Inches(3.2), Inches(0.5),
            font_size=13, color=GRIS_OSCURO)

# Derecha: tablas
tablas = [
    ("alertas_riesgo", "1,825 filas", "fecha, estacion_id, nivel_riesgo, ihb, itt, ieh, alerta_helada, alerta_tizon"),
    ("metricas_clima", "1,964 filas", "fecha, estacion_id, t_max, t_min, t_media, prec, humedad, et0"),
    ("predicciones_rendimiento", "5 filas", "variedad, rendimiento_predicho, rendimiento_historico, modelo"),
]
for i, (tabla, filas, cols) in enumerate(tablas):
    top = Inches(1.25 + i * 1.9)
    add_rect(sl, Inches(6.0), top, Inches(7.1), Inches(1.8), fill_color=BLANCO,
             line_color=AZUL, line_width=1.0)
    add_rect(sl, Inches(6.0), top, Inches(7.1), Inches(0.45), fill_color=RGBColor(0xE3,0xF2,0xFD))
    add_textbox(sl, "TABLE: " + tabla, Inches(6.1), top + Inches(0.05), Inches(4.5), Inches(0.35),
                font_size=13, bold=True, color=AZUL)
    add_textbox(sl, filas, Inches(11.2), top + Inches(0.05), Inches(1.7), Inches(0.35),
                font_size=13, bold=True, color=VERDE_OSCURO, align=PP_ALIGN.RIGHT)
    add_textbox(sl, "Columnas: " + cols, Inches(6.1), top + Inches(0.5), Inches(6.8), Inches(1.15),
                font_size=11, color=GRIS_OSCURO)

add_textbox(sl,
    "✅ Verificación real en EC2:\n"
    "cqlsh 54.159.9.135 → SELECT COUNT(*) = 1825 / 1964 / 5",
    Inches(6.0), Inches(7.0), Inches(7.1), Inches(0.4),
    font_size=12, color=VERDE_OSCURO, italic=True)

# ════════════════════════════════════════════════════════════
# SLIDE 10 – DASHBOARD STREAMLIT
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=GRIS_CLARO)
header_bar(sl, "Dashboard AgroSmart Andino",
           "Streamlit 1.x — 11 componentes visuales — Tiempo real")
footer_bar(sl); slide_number(sl, 10)

# Intentar insertar screenshot si existe
import os
screenshot_path = "d:/Documents/GitHub/Maestria_IA_caso_practico/Big_data/trabajo_final/agrosmart_andino/docs/dashboard_full.png"
if os.path.exists(screenshot_path):
    sl.shapes.add_picture(screenshot_path,
                          Inches(0.3), Inches(1.2), Inches(7.0), Inches(5.6))
    left_start = Inches(7.55)
    box_w = Inches(5.5)
else:
    left_start = Inches(0.3)
    box_w = Inches(12.7)
    add_textbox(sl, "[Screenshot dashboard_full.png]",
                Inches(0.3), Inches(1.2), Inches(7.0), Inches(5.6),
                font_size=14, color=GRIS_OSCURO, align=PP_ALIGN.CENTER)

# Componentes del dashboard
add_textbox(sl, "Componentes del Dashboard:",
            left_start, Inches(1.2), box_w, Inches(0.4),
            font_size=14, bold=True, color=VERDE_OSCURO)
componentes = [
    "🚦 Semáforo de riesgo (BAJO/MEDIO/ALTO/CRÍTICO)",
    "📊 KPIs: altitud, humedad, T.Media, T.Mínima",
    "⚡ Alertas activas: HELADA con timestamp",
    "🌡 Métricas climáticas: Prec 33.7mm, ET₀ 3.20",
    "📈 Serie temporal temperatura (Plotly)",
    "💧 Evolución precipitaciones y ET₀",
    "🍂 Índice de tizón tardío por semana",
    "🗺 Mapa de riesgo por estación (Plotly scatter)",
    "🔮 Gauge rendimiento predicho GBT",
    "📋 Tabla recomendaciones de riego",
    "🌾 Barras comparativas por variedad",
]
add_bullet_box(sl, componentes, left_start, Inches(1.7), box_w, Inches(5.15),
               font_size=12, bullet_color=GRIS_OSCURO, bullet_char="")

# ════════════════════════════════════════════════════════════
# SLIDE 11 – RESULTADOS
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=VERDE_OSCURO)
header_bar(sl, "Resultados Obtenidos",
           "Ejecución real AWS EMR | Datos verificados en EC2 Cassandra",
           bar_color=VERDE_MEDIO, title_color=BLANCO, sub_color=AMARILLO)
footer_bar(sl); slide_number(sl, 11)

metricas_grid = [
    (AMARILLO, "R² = 0.6764", "GBT Regressor\n(rendimiento t/ha)\nejecución real EMR"),
    (AMARILLO, "RMSE 0.97 t/ha", "Error predicción\nGBT sobre datos\nde prueba EMR"),
    (VERDE_CLARO, "72,382", "Registros procesados\npor Spark\n(Raw → Gold)"),
    (VERDE_CLARO, "3,794", "Filas cargadas\nen Cassandra\n(3 tablas NoSQL)"),
    (AZUL, "9.2 t/ha", "Rendimiento promedio\npredicho campaña\n2025-2026"),
    (NARANJA, "38 alertas HELADA", "Eventos críticos\ndetectados en\nhistórico 2015-2024"),
]
for i, (col, val, desc) in enumerate(metricas_grid):
    c = i % 3
    r = i // 3
    lft = Inches(0.35 + c * 4.3)
    top = Inches(1.2 + r * 2.6)
    add_rect(sl, lft, top, Inches(4.1), Inches(2.4), fill_color=col)
    add_textbox(sl, val, lft + Inches(0.15), top + Inches(0.2), Inches(3.8), Inches(0.95),
                font_size=26, bold=True, color=VERDE_OSCURO if col == AMARILLO else BLANCO,
                align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, lft + Inches(0.15), top + Inches(1.1), Inches(3.8), Inches(1.15),
                font_size=13, color=VERDE_OSCURO if col == AMARILLO else BLANCO,
                align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 12 – VIABILIDAD
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=GRIS_CLARO)
header_bar(sl, "Viabilidad, Escalabilidad e Impacto",
           "Expansión futura y posibilidades de crecimiento")
footer_bar(sl); slide_number(sl, 12)

columnas = [
    (VERDE_OSCURO, "🚀 Escenarios de uso",
     ["Agricultores de Huangascar con alertas SMS/WhatsApp",
      "Municipalidad de Yauyos para planificación de campaña",
      "SENAMHI para validación de datos climáticos locales",
      "MIDAGRI para estadísticas de rendimiento en tiempo real",
      "Universidades para investigación agroclimática"]),
    (AZUL, "⚡ Escalabilidad",
     ["EMR auto-scaling: 1→N nodos worker bajo demanda",
      "Cassandra replication factor expandible a 3+ nodos",
      "NASA POWER cubre todo el territorio peruano (API REST)",
      "Modelo GBT re-entrenable con nuevos datos de campaña",
      "Streamlit Cloud: deploy público sin infraestructura"]),
    (NARANJA, "⚠ Limitaciones actuales",
     ["IoT sintético (no sensores físicos desplegados)",
      "Cassandra en nodo único (RF=1, sin HA)",
      "Dashboard público en agrosmart-andino.streamlit.app",
      "Modelo entrenado con 1,500 registros históricos",
      "Sin integración con sistemas SENAMHI en tiempo real"]),
    (VERDE_CLARO, "🌱 Expansión futura",
     ["Integrar sensores IoT reales (bajo costo: ESP32+DHT22)",
      "Dashboard activo en agrosmart-andino.streamlit.app",
      "API REST para consumo desde apps móviles",
      "Modelo con datos satelitales (Sentinel-2 NDVI)",
      "Expansión a otras regiones andinas del Perú"]),
]
for i, (col, titulo, items) in enumerate(columnas):
    lft = Inches(0.3 + i * 3.25)
    add_rect(sl, lft, Inches(1.2), Inches(3.1), Inches(5.7), fill_color=BLANCO,
             line_color=col, line_width=1.5)
    add_rect(sl, lft, Inches(1.2), Inches(3.1), Inches(0.5), fill_color=col)
    add_textbox(sl, titulo, lft + Inches(0.05), Inches(1.25), Inches(3.0), Inches(0.4),
                font_size=13, bold=True, color=BLANCO)
    add_bullet_box(sl, items, lft + Inches(0.05), Inches(1.75), Inches(3.0), Inches(5.0),
                   font_size=11.5, bullet_color=GRIS_OSCURO)

# ════════════════════════════════════════════════════════════
# SLIDE 13 – CONCLUSIONES
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=VERDE_OSCURO)
header_bar(sl, "Conclusiones",
           "Experiencias del equipo y aprendizajes",
           bar_color=VERDE_MEDIO, title_color=BLANCO, sub_color=AMARILLO)
footer_bar(sl); slide_number(sl, 13)

conclusiones = [
    ("✅", "Pipeline Big Data completo",
     "Se diseñó e implementó un pipeline de 5 capas sobre AWS:\n"
     "NASA POWER/IoT → S3 → Spark EMR → Cassandra → Streamlit,\n"
     "procesando 72,382 registros en ejecución real."),
    ("✅", "Machine Learning con impacto real",
     "El GBT Regressor obtuvo R²=0.6764 y RMSE=0.97 t/ha en ejecución\n"
     "real sobre EMR 6.15, prediciendo rendimientos por variedad\n"
     "con valor operativo para agricultores andinos."),
    ("✅", "Alerta temprana funcional",
     "El sistema detectó 38 eventos HELADA críticos en el histórico\n"
     "2015-2024, demostrando la viabilidad del sistema de alertas\n"
     "para reducir pérdidas del 38% observadas en 2022."),
    ("✅", "Tecnologías Big Data integradas",
     "Uso efectivo de Apache Spark, Cassandra 4.1, AWS S3/EMR,\n"
     "Streamlit y PySpark MLlib en un contexto de aplicación real,\n"
     "validando la arquitectura propuesta en producción."),
    ("📚", "Aprendizajes del equipo",
     "La combinación de datos heterogéneos (clima + IoT + estadísticas)\n"
     "requiere ETL cuidadoso. El valor real no está en el algoritmo\n"
     "sino en la calidad del pipeline y la utilidad para el usuario final."),
]
for i, (icon, tit, desc) in enumerate(conclusiones):
    top = Inches(1.2 + i * 1.13)
    add_rect(sl, Inches(0.3), top, Inches(0.65), Inches(0.9), fill_color=AMARILLO)
    add_textbox(sl, icon, Inches(0.3), top + Inches(0.15), Inches(0.65), Inches(0.6),
                font_size=18, align=PP_ALIGN.CENTER, color=VERDE_OSCURO)
    add_textbox(sl, tit + ":", Inches(1.05), top + Inches(0.03), Inches(11.5), Inches(0.38),
                font_size=14, bold=True, color=AMARILLO)
    add_textbox(sl, desc, Inches(1.05), top + Inches(0.42), Inches(11.5), Inches(0.65),
                font_size=12, color=BLANCO)

# ════════════════════════════════════════════════════════════
# SLIDE 14 – CIERRE / Q&A
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill_color=VERDE_OSCURO)
add_rect(sl, 0, Inches(3.5), SLIDE_W, Inches(4.0), fill_color=VERDE_MEDIO)

add_textbox(sl, "🌱 AgroSmart Andino",
            Inches(1), Inches(0.5), Inches(11), Inches(1.2),
            font_size=44, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(1.5), Inches(1.75), Inches(10.3), Pt(3), fill_color=AMARILLO)
add_textbox(sl, "¡Gracias por su atención!",
            Inches(1), Inches(1.9), Inches(11), Inches(0.8),
            font_size=32, bold=True, color=AMARILLO, align=PP_ALIGN.CENTER)

add_textbox(sl, "Preguntas y Respuestas",
            Inches(1), Inches(3.7), Inches(11), Inches(0.7),
            font_size=28, bold=False, color=BLANCO, align=PP_ALIGN.CENTER)

# Info del repositorio
add_textbox(sl,
    "📂 Repositorio: github.com/wilmerjelko/Parcial_Caso_Practico_1\n"
    "📧 Grupo 2 | Maestría IA | UNI | Curso Big Data | Mayo 2026\n"
    "👩‍🏫 Docente: Mg. Rosa Virginia Encinas Quille",
    Inches(1.5), Inches(4.6), Inches(10), Inches(1.8),
    font_size=14, color=VERDE_FONDO, align=PP_ALIGN.CENTER)

# Guardar
output_path = "d:/Documents/GitHub/Maestria_IA_caso_practico/Big_data/trabajo_final/AgroSmart_Andino_Presentacion.pptx"
prs.save(output_path)
print(f"✅ Presentación guardada en:\n{output_path}")
print(f"   Total de slides: {len(prs.slides)}")
